import hashlib
import io
import logging
import re
import time
import uuid

logger = logging.getLogger(__name__)

import fitz
import streamlit as st
from google.genai import types
from PIL import Image, UnidentifiedImageError
from pptx import Presentation

MAX_IMAGE_EDGE = 1024
MAX_UPLOAD_MB = 50                          # hard cap: files larger than this are rejected
MAX_UPLOAD_BYTES = MAX_UPLOAD_MB * 1024 * 1024
IMAGE_ANALYSIS_PROMPT = "This is a computer science slide. Transcribe the code and explain any diagrams."
GEMINI_MODEL = "gemini-2.5-flash"

# Strict extension whitelist — nothing outside this set is processed
_ALLOWED_EXTENSIONS: frozenset[str] = frozenset({"pdf", "pptx", "jpg", "jpeg", "png", "txt"})

# Any extension matching these patterns is always rejected regardless of MIME
_BLOCKED_EXTENSIONS: frozenset[str] = frozenset({
    "db", "sqlite", "sqlite3", "sql",          # database files
    "py", "pyc", "pyw",                         # Python scripts
    "sh", "bash", "zsh", "bat", "cmd", "ps1",  # shell scripts
    "exe", "dll", "so", "dylib",               # executables / libraries
    "js", "ts", "php", "rb", "pl",             # other code
})


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def clean_to_markdown(raw_text: str) -> str:
    if not raw_text:
        return ""
    text = raw_text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            lines.append("")
            continue
        is_heading = (
            len(stripped) <= 90
            and len(stripped.split()) <= 12
            and not stripped.endswith((".", ",", ";"))
            and re.match(r"^(chapter|section|module|week|slide|\d+(\.\d+)*)\b", stripped, re.I)
        )
        lines.append(f"## {stripped}" if is_heading else stripped)
    return re.sub(r"\n{3,}", "\n\n", "\n".join(lines)).strip()


def parse_json_response(text: str) -> dict:
    import json
    cleaned = text.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    return json.loads(cleaned)


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def image_mime_from_ext(ext: str) -> str:
    return "image/jpeg" if ext.lower().lstrip(".") in {"jpg", "jpeg"} else "image/png"


def validate_image(image_bytes: bytes, mime_type: str) -> tuple[bytes, str]:
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            image.load()
            image.thumbnail((MAX_IMAGE_EDGE, MAX_IMAGE_EDGE))
            if image.mode not in {"RGB", "RGBA"}:
                image = image.convert("RGB")
            output = io.BytesIO()
            if image.mode == "RGBA" and mime_type == "image/png":
                image.save(output, format="PNG", optimize=True)
                return output.getvalue(), "image/png"
            image = image.convert("RGB")
            image.save(output, format="JPEG", quality=85, optimize=True)
            return output.getvalue(), "image/jpeg"
    except (UnidentifiedImageError, OSError, ValueError) as exc:
        raise ValueError("Visual analysis failed for this slide, but text was processed successfully.") from exc


# ---------------------------------------------------------------------------
# File extractors
# ---------------------------------------------------------------------------

def extract_pdf_text(file_bytes: bytes) -> tuple[str, int, str]:
    """Returns (text, page_count, warning)."""
    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            return "\n\n".join(page.get_text("text") for page in doc), doc.page_count, ""
    except Exception as exc:
        logger.error("PDF extraction failed: %s", exc, exc_info=True)
        return "", 0, "Unable to read this PDF — the file may be corrupted or password-protected."


def extract_pptx(file_bytes: bytes) -> tuple[str, list[dict], int, str]:
    """Returns (text, images, slide_count, warning)."""
    try:
        deck = Presentation(io.BytesIO(file_bytes))
        slides, images = [], []
        for slide_number, slide in enumerate(deck.slides, start=1):
            slide_text = [f"## Slide {slide_number}"]
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                if hasattr(shape, "image"):
                    img = shape.image
                    try:
                        image_bytes, mime_type = validate_image(img.blob, image_mime_from_ext(img.ext))
                        images.append({
                            "label": f"Slide {slide_number} image {shape_index}",
                            "bytes": image_bytes,
                            "mime_type": mime_type,
                        })
                    except ValueError:
                        continue
            slides.append("\n\n".join(slide_text))
        return "\n\n".join(slides), images, len(deck.slides), ""
    except Exception as exc:
        logger.error("PPTX extraction failed: %s", exc, exc_info=True)
        return "", [], 0, "Unable to read this PPTX — the file may be corrupted or use an unsupported format."


def analyze_image(api_key: str, image_bytes: bytes, mime_type: str) -> tuple[str, bool]:
    try:
        validated_bytes, validated_mime = validate_image(image_bytes, mime_type)
    except ValueError:
        return "Visual analysis failed for this slide, but text was processed successfully.", False
    if not api_key:
        return "", True
    try:
        from utils.gemini import get_gemini_client
        client = get_gemini_client(api_key)
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[IMAGE_ANALYSIS_PROMPT,
                      types.Part.from_bytes(data=validated_bytes, mime_type=validated_mime)],
        )
        return response.text or "", True
    except Exception:
        return "Visual analysis failed for this slide, but text was processed successfully.", False


# ---------------------------------------------------------------------------
# Main parse entry point
# ---------------------------------------------------------------------------

def parse_uploaded_file(uploaded_file, api_key: str) -> tuple[dict | None, int, list[str]]:
    """Parse one uploaded file. Returns (file_item_or_None, unit_count, warnings).

    Returns None as the first element when the file should be skipped entirely
    (oversized, empty content with no images, etc.).
    """
    from utils.metrics import report_parse_metrics
    parse_start = time.perf_counter()
    file_bytes = uploaded_file.getvalue()

    # ── Filename sanitisation — strip path traversal and non-safe chars ──────
    raw_name = uploaded_file.name or "upload"
    # Keep only the basename, drop any directory component
    safe_name = re.sub(r"[^\w.\- ]", "_", raw_name.split("/")[-1].split("\\")[-1])
    safe_name = safe_name.strip(". ") or "upload"
    file_name = safe_name

    file_size_bytes = len(file_bytes)

    # ── Extension guard ──────────────────────────────────────────────────────
    # Extract the extension after the last dot; reject double-extension tricks
    # like "report.pdf.py" by checking the true final extension.
    parts = file_name.rsplit(".", 1)
    file_ext = parts[-1].lower() if len(parts) == 2 else ""

    if file_ext in _BLOCKED_EXTENSIONS:
        logger.warning(
            "Blocked upload with dangerous extension '.%s' (original name: '%s')",
            file_ext, raw_name,
        )
        return None, 0, [
            f"'{raw_name}' was rejected: files with the .{file_ext} extension are not allowed."
        ]

    if file_ext not in _ALLOWED_EXTENSIONS:
        logger.warning(
            "Blocked upload with unknown extension '.%s' (original name: '%s')",
            file_ext, raw_name,
        )
        return None, 0, [
            f"'{raw_name}' is not a supported file type. "
            f"Allowed: {', '.join(sorted(_ALLOWED_EXTENSIONS))}."
        ]

    # ── Size guard ──────────────────────────────────────────────────────────
    if file_size_bytes > MAX_UPLOAD_BYTES:
        size_mb = file_size_bytes / (1024 * 1024)
        return None, 0, [
            f"'{file_name}' is {size_mb:.1f} MB — files larger than {MAX_UPLOAD_MB} MB are not supported. "
            "Split the document and re-upload."
        ]

    file_type = file_ext  # already extracted and validated above
    file_hash = hashlib.sha256(file_bytes).hexdigest()
    warnings, images = [], []
    indexed_units = 1
    raw_text = ""

    if file_type == "pdf":
        raw_text, indexed_units, warn = extract_pdf_text(file_bytes)
        if warn:
            warnings.append(warn)
    elif file_type == "pptx":
        raw_text, images, indexed_units, warn = extract_pptx(file_bytes)
        if warn:
            warnings.append(warn)
    elif file_type in {"jpg", "jpeg", "png"}:
        mime_type = image_mime_from_ext(file_type)
        visual_text, ok = analyze_image(api_key, file_bytes, mime_type)
        raw_text = f"## Visual Slide: {file_name}\n\n{visual_text}" if visual_text and ok else ""
        if ok:
            try:
                valid_bytes, valid_mime = validate_image(file_bytes, mime_type)
                images = [{"label": file_name, "bytes": valid_bytes, "mime_type": valid_mime}]
            except ValueError:
                pass
        else:
            warnings.append("Visual analysis failed for this slide, but text was processed successfully.")

    parse_end = time.perf_counter()
    cleaned_content = clean_to_markdown(raw_text)

    # ── Empty-content guard ────────────────────────────────────────────────
    # An image uploaded without an API key yields empty text and no usable
    # images. Skip it rather than inflating file/slide counts with dead weight.
    if not cleaned_content and not images:
        return None, 0, warnings + [
            f"'{file_name}' produced no content. "
            "For image files, a Gemini API key is required to extract text."
        ]

    report_parse_metrics(
        file_name=file_name,
        file_type=file_type,
        file_size_kb=file_size_bytes / 1024,
        pages_or_slides=indexed_units,
        raw_chars=len(raw_text),
        cleaned_chars=len(cleaned_content),
        images_found=len(images),
        parse_time_s=round(parse_end - parse_start, 2),
    )

    return (
        {"id": f"{file_name}:{file_hash}", "name": file_name, "type": file_type,
         "size": file_size_bytes, "hash": file_hash, "content": cleaned_content, "images": images},
        indexed_units,
        warnings,
    )


# ---------------------------------------------------------------------------
# Workspace helpers
# ---------------------------------------------------------------------------

def blank_workspace() -> dict:
    return {
        "id": uuid.uuid4().hex[:8],
        "files": [],
        "processed_text": "",
        "quiz_history": [],
        "generated_notes": "",
        "weak_area_report": "",
        # stats are always recomputed from files — never mutated in place
        "stats": {"pages": 0, "slides": 0, "chapters": 0},
        "visual_warnings": [],
        "quiz_attempt_counter": 0,
    }


def _recompute_stats(workspace: dict) -> None:
    """Single source of truth: derive all stats from workspace['files'] each time.

    pages  = sum of page/slide counts stored on each file item
    slides = same as pages (kept for display label compatibility)
    chapters = count of ## headings across all cleaned content
    """
    pages, chapters = 0, 0
    for f in workspace["files"]:
        pages += f.get("unit_count", 1)
        chapters += len(re.findall(r"^##\s+", f.get("content", ""), flags=re.M))
    workspace["stats"]["pages"] = pages
    workspace["stats"]["slides"] = pages      # alias kept so existing display code works
    workspace["stats"]["chapters"] = chapters


def refresh_processed_text(workspace: dict) -> None:
    """Rebuild processed_text and recompute all stats from scratch."""
    chunks = []
    for file_item in workspace["files"]:
        if file_item["content"]:
            chunks.append(f"# Source: {file_item['name']}\n\n{file_item['content']}")
    workspace["processed_text"] = "\n\n---\n\n".join(chunks).strip()
    _recompute_stats(workspace)


def add_textbook_content(workspace: dict, text: str, subject: str) -> bool:
    """Add pasted text as a file item. Returns True if new content was added."""
    cleaned = clean_to_markdown(text)
    if not cleaned:
        return False
    text_hash = hashlib.sha256(cleaned.encode("utf-8")).hexdigest()
    source_id = f"pasted-text:{text_hash}"
    if any(f["id"] == source_id for f in workspace["files"]):
        return False
    workspace["files"].append({
        "id": source_id,
        "name": f"{subject} pasted content",
        "type": "text",
        "size": len(cleaned),
        "hash": text_hash,
        "content": f"## Pasted Textbook Content\n\n{cleaned}",
        "images": [],
        "unit_count": max(1, len(re.findall(r"^##\s+", cleaned, flags=re.M))),
    })
    return True


def _persist_chunks(file_item: dict, workspace_id: str) -> None:
    """Chunk one file item and write rows to material_chunks.

    Only called for files that are genuinely new (not already in known_ids),
    so we never produce duplicate chunk rows for the same source file.
    Errors are logged and swallowed — a chunking failure must never block
    the main upload flow.
    """
    try:
        from utils.chunking import chunk_document
        from utils.persistence import SessionLocal, MaterialChunk

        chunks = chunk_document(
            content_text=file_item.get("content", ""),
            source_file_id=file_item["id"],
            workspace_id=workspace_id,
            source_name=file_item["name"],
        )
        if not chunks:
            return

        db = SessionLocal()
        try:
            for c in chunks:
                db.add(MaterialChunk(**c))
            db.commit()
        except Exception:
            db.rollback()
            logger.error(
                "Failed to persist chunks for '%s'", file_item["name"], exc_info=True
            )
        finally:
            db.close()
    except Exception:
        logger.error(
            "_persist_chunks crashed for '%s'", file_item.get("name"), exc_info=True
        )


def index_materials(uploaded_files, pasted_text: str, workspace: dict, subject: str, api_key: str) -> None:
    from utils.metrics import log_metric
    any_new = False
    warnings = []
    known_ids = {f["id"] for f in workspace["files"]}

    # workspace_id is needed for chunk foreign keys; fall back gracefully if absent.
    workspace_id: str = workspace.get("id", "")

    image_uploads = [
        item for item in uploaded_files or []
        if item.name.rsplit(".", 1)[-1].lower() in {"jpg", "jpeg", "png"}
    ]
    if image_uploads and not api_key:
        st.warning("⚙ A Gemini API key is required to analyze image slides.")

    for uploaded_file in uploaded_files or []:
        file_item, unit_count, file_warnings = parse_uploaded_file(uploaded_file, api_key)
        warnings.extend(file_warnings)
        if file_item is None:
            # Skipped by size cap or empty-content guard — warning already added
            continue
        if file_item["id"] not in known_ids:
            file_item["unit_count"] = unit_count   # store for stats recompute
            workspace["files"].append(file_item)
            known_ids.add(file_item["id"])
            any_new = True
            if workspace_id:
                _persist_chunks(file_item, workspace_id)

    if add_textbook_content(workspace, pasted_text, subject):
        any_new = True

    # Recompute everything from scratch — single source of truth
    refresh_processed_text(workspace)

    if any_new:
        st.session_state["is_dirty"] = True
        st.toast(
            f"Indexed! {workspace['stats']['slides']} pages/slides, "
            f"{workspace['stats']['chapters']} sections across {len(workspace['files'])} source(s)."
        )
    elif uploaded_files or pasted_text.strip():
        st.caption("Workspace already has this material indexed.")

    for warning in sorted(set(warnings)):
        st.warning(warning)
