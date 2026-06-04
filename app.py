import hashlib
import io
import json
import os
import re
from textwrap import dedent

import fitz
import streamlit as st
from google import genai
from google.genai import types
from PIL import Image, UnidentifiedImageError
from pptx import Presentation


APP_TITLE = "SunDevil AI"
GEMINI_MODEL = "models/gemini-3.5-flash"
SUPPORTED_UPLOADS = ["pdf", "pptx", "jpg", "jpeg", "png"]
MAX_IMAGE_EDGE = 1024
IMAGE_ANALYSIS_PROMPT = "This is a computer science slide. Transcribe the code and explain any diagrams."


def blank_workspace() -> dict:
    return {
        "files": [],
        "processed_text": "",
        "quiz_history": [],
        "generated_notes": "",
        "weak_area_report": "",
        "stats": {"slides": 0, "chapters": 0},
        "visual_warnings": [],
    }


def initialize_state() -> None:
    st.session_state.setdefault("workspaces", {"CSE 230": blank_workspace()})
    st.session_state.setdefault("active_workspace", next(iter(st.session_state["workspaces"])))


def active_workspace() -> dict:
    return st.session_state["workspaces"][st.session_state["active_workspace"]]


def apply_theme() -> None:
    st.markdown(
        """
        <style>
        :root {
            --bg: #FCF9F1;
            --sidebar: #F4F1E8;
            --panel: #FFFDF7;
            --ink: #2D2D2D;
            --muted: #6F6A60;
            --line: #E7DECF;
            --rose: #E8A0BF;
            --gold: #D4AF37;
            --gold-dark: #B8942F;
        }

        .stApp {
            background: var(--bg);
            color: var(--ink);
        }

        .main .block-container {
            max-width: 1360px;
            padding-top: 2rem;
        }

        [data-testid="stSidebar"] {
            background: var(--sidebar);
            border-right: 1px solid var(--line);
        }

        [data-testid="stSidebar"] * {
            color: var(--ink);
            letter-spacing: 0;
        }

        [data-testid="stSidebar"] .stButton > button {
            background: #8C1D40 !important;
            border: 1px solid #8C1D40 !important;
            color: #FFFFFF !important;
            border-radius: 999px !important;
            font-weight: 700 !important;
        }

        [data-testid="stSidebar"] .stButton > button *,
        [data-testid="stSidebar"] .stButton > button p,
        [data-testid="stSidebar"] .stButton > button span {
            color: #FFFFFF !important;
        }

        [data-testid="stSidebar"] .stButton > button:hover,
        [data-testid="stSidebar"] .stButton > button:focus {
            background: #741634 !important;
            border-color: #741634 !important;
            color: #FFFFFF !important;
        }

        h1, h2, h3, p, label, span {
            color: var(--ink);
            letter-spacing: 0;
        }

        div[data-testid="stTabs"] button[role="tab"] {
            border-radius: 999px;
            padding: 10px 16px;
            color: var(--ink);
        }

        div[data-testid="stTabs"] button[aria-selected="true"] {
            background: #FFFFFF;
            border: 1px solid var(--line);
        }

        .stButton > button,
        .stDownloadButton > button {
            border-radius: 999px;
            border: 1px solid var(--line);
            color: var(--ink) !important;
            font-weight: 650;
        }

        .stButton > button[kind="primary"],
        .stDownloadButton > button {
            background: var(--gold);
            border-color: var(--gold);
            color: #2D2D2D !important;
            font-weight: 700;
        }

        .stButton > button *,
        .stButton > button p,
        .stButton > button span,
        .stDownloadButton > button *,
        .stDownloadButton > button p,
        .stDownloadButton > button span {
            color: inherit !important;
        }

        .stButton > button[kind="primary"]:hover,
        .stDownloadButton > button:hover {
            background: var(--gold-dark);
            border-color: var(--gold-dark);
            color: #2D2D2D !important;
        }

        input,
        textarea,
        div[data-baseweb="input"] input,
        div[data-baseweb="textarea"] textarea,
        [data-testid="stTextInput"] input,
        [data-testid="stTextArea"] textarea {
            color: #2D2D2D !important;
            caret-color: #2D2D2D !important;
            background: #FFFFFF !important;
        }

        input::placeholder,
        textarea::placeholder,
        [data-testid="stTextInput"] input::placeholder,
        [data-testid="stTextArea"] textarea::placeholder {
            color: #6F6A60 !important;
            opacity: 1 !important;
        }

        .stFileUploader,
        div[data-testid="stMetric"],
        div[data-testid="stTextArea"],
        div[data-testid="stExpander"] {
            background: var(--panel);
            border-radius: 14px;
        }

        div[data-testid="stTextArea"] textarea {
            border: 2px solid var(--rose);
            border-radius: 14px;
            background: #FFFFFF;
            color: var(--ink);
        }

        div[data-testid="stAlert"] {
            border-radius: 14px;
            color: var(--ink);
        }

        .workspace-card {
            background: var(--panel);
            border: 1px solid var(--line);
            border-radius: 16px;
            padding: 18px 20px;
            margin: 10px 0 18px 0;
        }

        .workspace-card small {
            color: var(--muted);
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def clean_to_markdown(raw_text: str) -> str:
    if not raw_text:
        return ""

    text = raw_text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)

    lines = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            lines.append("")
            continue

        is_heading = (
            len(stripped) <= 90
            and len(stripped.split()) <= 12
            and not stripped.endswith((".", ",", ";"))
            and re.match(r"^(chapter|section|module|week|slide|\d+(\.\d+)*)\b", stripped, re.I)
        )
        lines.append(f"## {stripped}" if is_heading else stripped)

    return re.sub(r"\n{3,}", "\n\n", "\n".join(lines)).strip()


def extract_pdf_text(file_bytes: bytes) -> tuple[str, int]:
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        return "\n\n".join(page.get_text("text") for page in doc), doc.page_count


def image_mime_from_ext(ext: str) -> str:
    return "image/jpeg" if ext.lower().lstrip(".") in {"jpg", "jpeg"} else "image/png"


def validate_image(image_bytes: bytes, mime_type: str) -> tuple[bytes, str]:
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            image.load()
            image.thumbnail((MAX_IMAGE_EDGE, MAX_IMAGE_EDGE))

            if image.mode not in {"RGB", "RGBA"}:
                image = image.convert("RGB")

            output = io.BytesIO()
            if image.mode == "RGBA" and mime_type == "image/png":
                image.save(output, format="PNG", optimize=True)
                return output.getvalue(), "image/png"

            image = image.convert("RGB")
            image.save(output, format="JPEG", quality=85, optimize=True)
            return output.getvalue(), "image/jpeg"
    except (UnidentifiedImageError, OSError, ValueError) as exc:
        raise ValueError("Visual analysis failed for this slide, but text was processed successfully.") from exc


def extract_pptx(file_bytes: bytes) -> tuple[str, list[dict], int]:
    deck = Presentation(io.BytesIO(file_bytes))
    slides = []
    images = []

    for slide_number, slide in enumerate(deck.slides, start=1):
        slide_text = [f"## Slide {slide_number}"]
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            if hasattr(shape, "image"):
                image = shape.image
                try:
                    image_bytes, mime_type = validate_image(image.blob, image_mime_from_ext(image.ext))
                    images.append(
                        {
                            "label": f"Slide {slide_number} image {shape_index}",
                            "bytes": image_bytes,
                            "mime_type": mime_type,
                        }
                    )
                except ValueError:
                    continue
        slides.append("\n\n".join(slide_text))

    return "\n\n".join(slides), images, len(deck.slides)


def analyze_image(api_key: str, image_bytes: bytes, mime_type: str) -> tuple[str, bool]:
    try:
        validated_bytes, validated_mime = validate_image(image_bytes, mime_type)
    except ValueError:
        return "Visual analysis failed for this slide, but text was processed successfully.", False

    if not api_key:
        return "", True

    try:
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[
                IMAGE_ANALYSIS_PROMPT,
                types.Part.from_bytes(data=validated_bytes, mime_type=validated_mime),
            ],
        )
        return response.text or "", True
    except Exception:
        return "Visual analysis failed for this slide, but text was processed successfully.", False


def parse_uploaded_file(uploaded_file, api_key: str) -> tuple[dict, int, list[str]]:
    file_bytes = uploaded_file.getvalue()
    file_name = uploaded_file.name
    file_type = file_name.rsplit(".", 1)[-1].lower()
    file_hash = hashlib.sha256(file_bytes).hexdigest()
    warnings = []
    indexed_units = 1
    images = []

    if file_type == "pdf":
        raw_text, indexed_units = extract_pdf_text(file_bytes)
    elif file_type == "pptx":
        raw_text, images, indexed_units = extract_pptx(file_bytes)
    elif file_type in {"jpg", "jpeg", "png"}:
        mime_type = image_mime_from_ext(file_type)
        visual_text, ok = analyze_image(api_key, file_bytes, mime_type)
        raw_text = f"## Visual Slide: {file_name}\n\n{visual_text}" if visual_text and ok else ""
        if ok:
            try:
                valid_bytes, valid_mime = validate_image(file_bytes, mime_type)
                images = [{"label": file_name, "bytes": valid_bytes, "mime_type": valid_mime}]
            except ValueError:
                pass
        else:
            warnings.append("Visual analysis failed for this slide, but text was processed successfully.")
    else:
        raw_text = ""

    return (
        {
            "id": f"{file_name}:{file_hash}",
            "name": file_name,
            "type": file_type,
            "size": len(file_bytes),
            "hash": file_hash,
            "content": clean_to_markdown(raw_text),
            "images": images,
        },
        indexed_units,
        warnings,
    )


def refresh_processed_text(workspace: dict) -> None:
    chunks = []
    chapter_count = 0
    for file_item in workspace["files"]:
        if file_item["content"]:
            chunks.append(f"# Source: {file_item['name']}\n\n{file_item['content']}")
            chapter_count += len(re.findall(r"^##\s+", file_item["content"], flags=re.M))

    workspace["processed_text"] = "\n\n---\n\n".join(chunks).strip()
    workspace["stats"]["chapters"] = chapter_count


def add_textbook_content(workspace: dict, text: str, subject: str) -> int:
    cleaned = clean_to_markdown(text)
    if not cleaned:
        return 0

    text_hash = hashlib.sha256(cleaned.encode("utf-8")).hexdigest()
    source_id = f"pasted-text:{text_hash}"
    if any(file_item["id"] == source_id for file_item in workspace["files"]):
        return 0

    workspace["files"].append(
        {
            "id": source_id,
            "name": f"{subject} pasted content",
            "type": "text",
            "size": len(cleaned),
            "hash": text_hash,
            "content": f"## Pasted Textbook Content\n\n{cleaned}",
            "images": [],
        }
    )
    refresh_processed_text(workspace)
    return max(1, len(re.findall(r"^##\s+", cleaned, flags=re.M)))


def index_materials(uploaded_files, pasted_text: str, workspace: dict, subject: str, api_key: str) -> None:
    indexed_units = 0
    warnings = []
    known_ids = {file_item["id"] for file_item in workspace["files"]}

    image_uploads = [
        item for item in uploaded_files or [] if item.name.rsplit(".", 1)[-1].lower() in {"jpg", "jpeg", "png"}
    ]
    if image_uploads and not api_key:
        st.warning("⚙ Setup Required")

    for uploaded_file in uploaded_files or []:
        file_item, units, file_warnings = parse_uploaded_file(uploaded_file, api_key)
        if file_item["id"] not in known_ids:
            workspace["files"].append(file_item)
            known_ids.add(file_item["id"])
            indexed_units += units
        warnings.extend(file_warnings)

    indexed_units += add_textbook_content(workspace, pasted_text, subject)
    workspace["stats"]["slides"] += indexed_units
    refresh_processed_text(workspace)

    if indexed_units:
        st.toast(f"Workspace Loaded: {workspace['stats']['slides']} Slides, {workspace['stats']['chapters']} Chapters")
    elif uploaded_files or pasted_text.strip():
        st.caption("Workspace already has this material indexed.")

    for warning in sorted(set(warnings)):
        st.warning(warning)


def workspace_image_parts(workspace: dict) -> list[types.Part]:
    parts = []
    for file_item in workspace["files"]:
        for image in file_item["images"]:
            try:
                image_bytes, mime_type = validate_image(image["bytes"], image["mime_type"])
                parts.append(types.Part.from_bytes(data=image_bytes, mime_type=mime_type))
            except ValueError:
                workspace["visual_warnings"].append(
                    "Visual analysis failed for this slide, but text was processed successfully."
                )
    return parts


def call_gemini(api_key: str, prompt: str, workspace: dict) -> str:
    image_parts = workspace_image_parts(workspace)
    try:
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[prompt, *image_parts],
        )
        return response.text or ""
    except Exception as exc:
        if image_parts:
            workspace["visual_warnings"].append(
                "Visual analysis failed for this slide, but text was processed successfully."
            )
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(model=GEMINI_MODEL, contents=[prompt])
            return response.text or ""
        raise RuntimeError(f"Gemini request failed: {exc}") from exc


def guide_prompt(subject: str, workspace: dict, mode: str) -> str:
    visual_instruction = ""
    if workspace_image_parts(workspace):
        visual_instruction = (
            "Attached images may include computer science slides, code screenshots, diagrams, stack/heap layouts, "
            "or memory traces. Explain the visual material where it supports the topic."
        )

    mode_instruction = (
        "Go deep on theory, edge cases, implementation tradeoffs, and memory/performance implications."
        if mode == "Deep Dive"
        else "Prioritize high-yield exam facts, syntax patterns, common pitfalls, and fast recall."
    )

    return dedent(
        f"""
        You are a PhD Teaching Assistant for ASU Computer Science.
        Subject workspace: {subject}
        {mode_instruction}
        {visual_instruction}

        Generate a Physics Method study guide using ONLY this workspace's materials.
        For every major topic, output exactly:

        ## [Topic Name]
        **THE RULE**: Explain the core rule, why it works, and memory/performance implications.
        **THE GUIDED SOLVE**: Provide a code or logic block with line-by-line documentation.
        **THE CHALLENGE**: Provide a practice problem that tests transfer and edge cases.
        **[ANSWER]**: Provide the complete worked answer and reasoning.

        Workspace materials:
        {workspace["processed_text"] if workspace["processed_text"] else "No text was extracted. Use attached images."}
        """
    ).strip()


def quiz_prompt(subject: str, workspace: dict) -> str:
    return dedent(
        f"""
        You are a PhD Teaching Assistant for ASU Computer Science.
        Subject workspace: {subject}
        Generate exactly 5 multiple-choice quiz questions using ONLY this workspace's materials.
        Use attached visuals when relevant.

        Return strict JSON only. No Markdown fences.
        Format:
        {{
          "questions": [
            {{
              "question": "Question text",
              "choices": ["A", "B", "C", "D"],
              "answer_index": 0,
              "topic": "Topic name",
              "explanation": "Why the correct answer is right"
            }}
          ]
        }}

        Workspace materials:
        {workspace["processed_text"] if workspace["processed_text"] else "No text was extracted. Use attached images."}
        """
    ).strip()


def parse_json_response(text: str) -> dict:
    cleaned = text.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    return json.loads(cleaned)


def split_topics(markdown: str) -> list[dict[str, str]]:
    topics = []
    current_title = ""
    current_lines = []

    for line in markdown.splitlines():
        heading = re.match(r"^##\s+(.+?)\s*$", line)
        if heading:
            if current_title:
                topics.append(topic_from_lines(current_title, current_lines))
            current_title = heading.group(1).strip()
            current_lines = []
        elif current_title:
            current_lines.append(line)

    if current_title:
        topics.append(topic_from_lines(current_title, current_lines))

    return topics


def topic_from_lines(title: str, lines: list[str]) -> dict[str, str]:
    body = "\n".join(lines).strip()
    answer_match = re.search(r"\*\*\[ANSWER\]\*\*\s*:\s*", body)
    if not answer_match:
        return {"title": title, "body": body, "answer": "_No answer section returned._"}
    return {
        "title": title,
        "body": body[: answer_match.start()].strip(),
        "answer": body[answer_match.end() :].strip(),
    }


def render_guide(markdown: str) -> None:
    topics = split_topics(markdown)
    if not topics:
        st.markdown(markdown)
        return

    for topic in topics:
        st.subheader(topic["title"])
        st.markdown(topic["body"])
        with st.expander("Click to see answer"):
            st.markdown(topic["answer"])


def workspace_summary(workspace: dict) -> None:
    st.markdown(
        f"""
        <div class="workspace-card">
          <strong>Workspace Loaded:</strong> {workspace["stats"]["slides"]} Slides, {workspace["stats"]["chapters"]} Chapters
          <br><small>{len(workspace["files"])} indexed source(s)</small>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_ingest_tab(subject: str, workspace: dict, api_key: str) -> None:
    subjects = list(st.session_state["workspaces"].keys())
    selected_subject = st.selectbox(
        "Subject Workspace",
        subjects,
        index=subjects.index(subject),
        key="ingest_subject_selector",
    )
    if selected_subject != subject:
        st.session_state["active_workspace"] = selected_subject
        st.rerun()

    left, right = st.columns([1, 1], gap="large")
    with left:
        uploaded_files = st.file_uploader(
            "Upload source files",
            type=SUPPORTED_UPLOADS,
            accept_multiple_files=True,
            key=f"uploader_{subject}",
            help="Accepted: PDF, PPTX, JPG, PNG.",
        )
    with right:
        pasted_text = st.text_area(
            "Paste Text",
            height=240,
            placeholder="Paste Zybooks or textbook content for this subject...",
            key=f"textbook_{subject}",
        )

    if st.button("Index Materials", type="primary"):
        index_materials(uploaded_files, pasted_text, workspace, subject, api_key)

    workspace_summary(workspace)
    for warning in sorted(set(workspace["visual_warnings"])):
        st.warning(warning)


def render_study_tab(api_key: str, subject: str, workspace: dict, mode: str) -> None:
    if st.button("Generate Physics Method Guide", type="primary"):
        if not api_key:
            st.warning("⚙ Setup Required")
        elif not workspace["files"]:
            st.warning("Add material in the Ingest Material tab first.")
        else:
            with st.spinner("Building study guide..."):
                try:
                    workspace["generated_notes"] = call_gemini(api_key, guide_prompt(subject, workspace, mode), workspace)
                except Exception as exc:
                    st.error(str(exc))

    if workspace["generated_notes"]:
        st.download_button(
            "Download Study Guide",
            data=workspace["generated_notes"],
            file_name=f"{subject.lower().replace(' ', '_')}_study_guide.md",
            mime="text/markdown",
            type="primary",
        )
        render_guide(workspace["generated_notes"])
    else:
        st.caption("Generate a Physics Method guide after indexing material.")


def render_quiz_tab(api_key: str, subject: str, workspace: dict) -> None:
    quiz_key = f"quiz_{subject}"
    answer_key = f"answers_{subject}"
    st.session_state.setdefault(quiz_key, [])
    st.session_state.setdefault(answer_key, {})

    if st.button("Generate Quiz", type="primary"):
        if not api_key:
            st.warning("⚙ Setup Required")
        elif not workspace["files"]:
            st.warning("Add material in the Ingest Material tab first.")
        else:
            with st.spinner("Generating quiz..."):
                try:
                    response_text = call_gemini(api_key, quiz_prompt(subject, workspace), workspace)
                    st.session_state[quiz_key] = parse_json_response(response_text).get("questions", [])[:5]
                    st.session_state[answer_key] = {}
                except Exception as exc:
                    st.error(str(exc))

    quiz = st.session_state[quiz_key]
    if not quiz:
        st.caption("Generate a quiz from the active workspace.")
        return

    for index, question in enumerate(quiz):
        choices = question.get("choices", [])
        st.markdown(f"**Q{index + 1}. {question.get('question', '')}**")
        if not choices:
            st.warning("This generated question did not include choices.")
            continue
        selected = st.radio(
            "Choose one",
            choices,
            key=f"{subject}_answer_{index}",
            label_visibility="collapsed",
        )
        st.session_state[answer_key][str(index)] = choices.index(selected)

    if st.button("Submit Quiz", type="primary"):
        correct = 0
        missed = []
        for index, question in enumerate(quiz):
            if st.session_state[answer_key].get(str(index)) == question.get("answer_index"):
                correct += 1
            else:
                missed.append(question)

        score = round((correct / len(quiz)) * 100) if quiz else 0
        workspace["quiz_history"].append(
            {
                "score": score,
                "questions": quiz,
                "answers": dict(st.session_state[answer_key]),
                "missed_questions": missed,
            }
        )
        st.success(f"Quiz saved: {score}%")


def render_workspace_sidebar() -> tuple[str, str, str]:
    with st.sidebar:
        st.header("Workspace")
        new_subject = st.text_input("New Subject Name", placeholder="CSE 230, Physics, MAT 243")
        if st.button("Create Workspace", type="primary", use_container_width=True):
            subject = new_subject.strip()
            if subject:
                st.session_state["workspaces"].setdefault(subject, blank_workspace())
                st.session_state["active_workspace"] = subject
                st.rerun()

        subjects = list(st.session_state["workspaces"].keys())
        active = st.session_state.get("active_workspace", subjects[0])
        if active not in subjects:
            active = subjects[0]

        selected = st.radio("Switch Workspace", subjects, index=subjects.index(active), key="workspace_selector")
        st.session_state["active_workspace"] = selected

        if st.button("Delete Workspace", type="primary", use_container_width=True):
            if len(st.session_state["workspaces"]) == 1:
                st.warning("Create another workspace before deleting the last one.")
            else:
                del st.session_state["workspaces"][selected]
                st.session_state["active_workspace"] = next(iter(st.session_state["workspaces"]))
                st.rerun()

        st.divider()
        with st.expander("⚙ Settings"):
            api_key = st.text_input("Gemini Key", value=os.getenv("GEMINI_API_KEY", ""), type="password")
            st.caption(f"Model: `{GEMINI_MODEL}`")
        study_mode = st.radio("Study Mode", ["Deep Dive", "Cram Mode"])

    return selected, api_key, study_mode


def main() -> None:
    st.set_page_config(page_title=APP_TITLE, page_icon=":books:", layout="wide")
    initialize_state()
    apply_theme()

    subject, api_key, study_mode = render_workspace_sidebar()
    workspace = active_workspace()

    st.title(APP_TITLE)
    st.caption("Warm, focused study workspaces for ASU computer science courses.")
    st.subheader(subject)

    ingest_tab, guide_tab, quiz_tab = st.tabs(["Ingest Material", "Study Guide", "Interactive Quiz"])

    with ingest_tab:
        render_ingest_tab(subject, workspace, api_key)
    with guide_tab:
        render_study_tab(api_key, subject, workspace, study_mode)
    with quiz_tab:
        render_quiz_tab(api_key, subject, workspace)


if __name__ == "__main__":
    main()
